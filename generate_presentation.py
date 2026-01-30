"""
Generate PowerPoint presentation for LaptopAI project
Author: Mario Radoš
Course: Znanstveno programiranje
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    TITLE_COLOR = RGBColor(31, 78, 121)  # Dark blue
    ACCENT_COLOR = RGBColor(68, 114, 196)  # Blue
    TEXT_COLOR = RGBColor(64, 64, 64)  # Dark gray
    
    def add_title_slide(title, subtitle):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = TITLE_COLOR
        title_para.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = ACCENT_COLOR
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, content_items, layout_type="bullet"):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = TITLE_COLOR
        
        # Underline
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(2)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.4), Inches(5))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(10)
            if layout_type == "bullet":
                p.level = 0
        
        return slide
    
    def add_two_column_slide(title, left_items, right_items):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = TITLE_COLOR
        
        # Underline
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
        line.line.color.rgb = ACCENT_COLOR
        line.line.width = Pt(2)
        
        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(4), Inches(5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        for i, item in enumerate(left_items):
            if i > 0:
                left_frame.add_paragraph()
            p = left_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)
        
        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.7), Inches(4), Inches(5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        for i, item in enumerate(right_items):
            if i > 0:
                right_frame.add_paragraph()
            p = right_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(8)
        
        return slide
    
    # Slide 1: Title
    add_title_slide(
        "LaptopAI",
        "AI-powered analiza korisničkih recenzija laptopa"
    )
    
    # Slide 2: Author info
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    info_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    text_frame = info_box.text_frame
    text_frame.text = "Mario Radoš"
    p0 = text_frame.paragraphs[0]
    p0.font.size = Pt(36)
    p0.font.bold = True
    p0.font.color.rgb = TITLE_COLOR
    p0.alignment = PP_ALIGN.CENTER
    
    text_frame.add_paragraph()
    p1 = text_frame.paragraphs[1]
    p1.text = "Prirodoslovno-matematički fakultet u Splitu"
    p1.font.size = Pt(20)
    p1.font.color.rgb = TEXT_COLOR
    p1.alignment = PP_ALIGN.CENTER
    
    text_frame.add_paragraph()
    p2 = text_frame.paragraphs[2]
    p2.text = "Znanstveno programiranje"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_COLOR
    p2.alignment = PP_ALIGN.CENTER
    
    text_frame.add_paragraph()
    p3 = text_frame.paragraphs[3]
    p3.text = "2026"
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_COLOR
    p3.alignment = PP_ALIGN.CENTER
    
    # Slide 3: Problem Statement
    add_content_slide(
        "Problem",
        [
            "🔍 Previše rasutih recenzija – korisnici dijele iskustva na desetke subreddita",
            "",
            "⚖️ Nemoguće usporedbe – teško je usporediti iskustva različitih korisnika",
            "",
            "⏱️ Gubljen vremena – nitko ne želi čitati stotine postova i komentara",
            "",
            "❓ Konfuzne informacije – službene recenzije često nedovoljne ili pristrane"
        ]
    )
    
    # Slide 4: Solution
    add_content_slide(
        "Rješenje: LaptopAI Pipeline",
        [
            "1️⃣ Prikupljanje – pronalazi relevantne Reddit rasprave o laptopima",
            "",
            "2️⃣ Pohranjivanje – sprema ih u semantički pretraživu bazu znanja",
            "",
            "3️⃣ Analiza – koristi AI za ekstrakciju sentimenta, pros/cons",
            "",
            "4️⃣ Rezultat – isporučuje čistu, strukturiranu preporuku i usporedbu"
        ]
    )
    
    # Slide 5: Key Features
    add_content_slide(
        "Ključne funkcionalnosti",
        [
            "🔍 Reddit scraping – automatsko prikupljanje korisničkih recenzija",
            "",
            "🧠 Semantičko pretraživanje – embeddinzi omogućuju pronalaženje relevantnih informacija",
            "",
            "🤖 AI analiza sentimenta – Google Gemini ekstrahira pros/cons i ocjenjuje",
            "",
            "⚔️ Laptop Battle UI – web sučelje za usporedbu dva laptopa u realnom vremenu",
            "",
            "📁 Modularni dizajn – nezavisni scraper, vector store i LLM slojevi"
        ]
    )
    
    # Slide 6: Tech Stack
    add_two_column_slide(
        "Tehnologije",
        [
            "Backend:",
            "• Python 3.11+",
            "• FastAPI (REST API)",
            "• ChromaDB (vector baza)",
            "• SentenceTransformers",
            "• Google Gemini API",
            "",
            "Scraping:",
            "• BeautifulSoup4",
            "• Requests"
        ],
        [
            "Frontend:",
            "• React 18",
            "• Vite",
            "• TailwindCSS",
            "• Axios",
            "",
            "Deployment:",
            "• Uvicorn ASGI server",
            "• Node.js dev server"
        ]
    )
    
    # Slide 7: System Architecture
    add_content_slide(
        "Arhitektura sustava",
        [
            "Korisnik (Web UI)",
            "         ⬇️",
            "FastAPI Backend (/api/compare)",
            "         ⬇️",
            "Data Collection Layer (Reddit Scraping)",
            "         ⬇️",
            "Knowledge Storage Layer (ChromaDB + Embeddings)",
            "         ⬇️",
            "Analysis Layer (Google Gemini LLM)",
            "         ⬇️",
            "Output Layer (JSON Results + Cache)"
        ],
        layout_type="plain"
    )
    
    # Slide 8: Data Collection Layer
    add_content_slide(
        "1. Data Collection Layer",
        [
            "📂 scrape_query_reddit.py – pronalazi relevantne postove",
            "   • BeautifulSoup + Requests (HTML scraping)",
            "   • Reddit search: https://reddit.com/search/?q=laptop+name",
            "   • Ekstrahira /comments/ linkove",
            "",
            "📄 reddit_post_scrapper.py – ekstrahira sadržaj",
            "   • Parsira <h1> (naslov), <div> (body), <p> (komentari)",
            "   • Sprema u data/reddit/<laptop-slug>/*.json",
            "",
            "❓ Zašto HTML scraping?",
            "   • Reddit API (PRAW) odbio pristup",
            "   • HTML scraping radi bez OAuth autentifikacije"
        ]
    )
    
    # Slide 9: Knowledge Storage Layer
    add_content_slide(
        "2. Knowledge Storage Layer",
        [
            "🗄️ ChromaDB – lokalna vector baza podataka",
            "",
            "📐 Embeddinzi: SentenceTransformers (all-MiniLM-L6-v2)",
            "   • 384-dimenzionalni vektori",
            "   • Semantičko pretraživanje (ne samo keyword matching)",
            "",
            "💾 Što se pohranjuje:",
            "   • Tekst postova i komentara",
            "   • Metapodaci (URL, subreddit, timestamp)",
            "   • Vector embeddinzi",
            "",
            "✅ Zašto ChromaDB? Brzo, lokalno, skalabilno, RAG-friendly"
        ]
    )
    
    # Slide 10: Analysis Layer
    add_content_slide(
        "3. Analysis Layer (LLM)",
        [
            "🤖 Google Gemini 1.5 Pro",
            "",
            "Proces:",
            "1. Query ChromaDB s laptop imenom",
            "2. Dohvati top 10-20 relevantnih postova",
            "3. Konstruiraj prompt s kontekstom",
            "4. Šalji LLM-u (Retrieval Augmented Generation)",
            "5. Parsiraj JSON odgovor",
            "",
            "Output format:",
            "• sentiment_score (1-100)",
            "• pros/cons liste",
            "• key_themes",
            "• user_recommendation"
        ]
    )
    
    # Slide 11: API Layer
    add_content_slide(
        "4. FastAPI Backend",
        [
            "🔌 Endpoint: POST /api/compare",
            "",
            "Input:",
            '{"laptop1": "Lenovo Legion Y540", "laptop2": "Dell XPS 15"}',
            "",
            "Logika:",
            "1. Caching check (postoje li već JSONovi?)",
            "2. Paralelna analiza za oba laptopa",
            "3. Usporedba sentiment_score → određuje winnera",
            "4. Vraća JSON response",
            "",
            "⚡ Cache: <200ms umjesto ~10s (bez LLM poziva)"
        ]
    )
    
    # Slide 12: Data Flow
    add_content_slide(
        "Tok podataka (puni pipeline)",
        [
            "Korisnik: 'Lenovo Legion Y540'",
            "  ⬇️",
            "Reddit Search → ~50 relevantnih postova",
            "  ⬇️",
            "Post Scraper → title, body, komentari",
            "  ⬇️",
            "ChromaDB → embedira i sprema",
            "  ⬇️",
            "Gemini LLM → generira pros/cons, sentiment",
            "  ⬇️",
            "JSON output → analysis/lenovo_legion_y540.json",
            "  ⬇️",
            "API Response → frontend (React UI)"
        ],
        layout_type="plain"
    )
    
    # Slide 13: Example Output
    add_content_slide(
        "Primjer rezultata",
        [
            "Laptop: Lenovo Legion Y540",
            "Sentiment Score: 78/100",
            "",
            "Pros:",
            "  • Odličan omjer cijene i performansi",
            "  • Dobro hlađenje uz RTX 2060",
            "  • Kvalitetna tipkovnica",
            "",
            "Cons:",
            "  • Loša baterija (2-3 sata)",
            "  • Osrednji ekran (sRGB ~60%)",
            "  • Plastični build quality",
            "",
            "Preporuka: Gaming na budžetu ✅, profesionalna upotreba ❌"
        ]
    )
    
    # Slide 14: Design Decisions
    add_two_column_slide(
        "Dizajnerske odluke",
        [
            "ChromaDB vs SQL:",
            "✅ Semantičko pretraživanje",
            "✅ Automatski embeddinzi",
            "✅ Skalira s velikim tekstovima",
            "✅ RAG-friendly",
            "",
            "Modularni pipeline:",
            "✅ Testabilnost",
            "✅ Skalabilnost",
            "✅ Debugging",
            "✅ Reusability"
        ],
        [
            "Zašto JSON output?:",
            "✅ Frontend-friendly",
            "✅ Human-readable",
            "✅ Strukturirano",
            "✅ LLM-friendly",
            "",
            "Zašto .gitignore data/?:",
            "⚠️ Veliki fileovi (>100MB)",
            "⚠️ Često se mijenjaju",
            "✅ Reproducibilni",
            "⚠️ Copyright-protected"
        ]
    )
    
    # Slide 15: Limitations
    add_content_slide(
        "Ograničenja",
        [
            "🔴 Ovisnost o kvaliteti Reddit podataka",
            "   • Ako nema dovoljno postova → loša analiza",
            "   • Stari postovi (>2 god) → možda nije relevantno",
            "",
            "🔴 AI output varijabilnost",
            "   • LLM nije deterministički (različiti rezultati)",
            "   • Sentiment score nije egzaktna znanost",
            "",
            "🔴 Nema real-time scrapinga",
            "   • Pipeline se pokreće ručno",
            "   • Novi postovi ne prikupljaju se automatski",
            "",
            "🔴 Rate limiting",
            "   • Reddit može blokirati previše requesta"
        ]
    )
    
    # Slide 16: DEMO PLACEHOLDER
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    demo_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1.5))
    demo_frame = demo_box.text_frame
    demo_frame.text = "🎬 DEMO"
    demo_para = demo_frame.paragraphs[0]
    demo_para.font.size = Pt(72)
    demo_para.font.bold = True
    demo_para.font.color.rgb = ACCENT_COLOR
    demo_para.alignment = PP_ALIGN.CENTER
    
    demo_frame.add_paragraph()
    p1 = demo_frame.paragraphs[1]
    p1.text = "Live demonstracija aplikacije"
    p1.font.size = Pt(24)
    p1.font.color.rgb = TEXT_COLOR
    p1.alignment = PP_ALIGN.CENTER
    
    # Slide 17: Future Work
    add_content_slide(
        "Budući razvoj",
        [
            "🔮 Više izvora podataka",
            "   • YouTube transkripti, tech forumi, LaptopMedia",
            "",
            "🔮 Automatski scheduled scraping",
            "   • Cron job za dnevno osvježavanje",
            "   • Background tasks (Celery)",
            "",
            "🔮 Historijski tracking",
            "   • Praćenje promjena sentimenta kroz vrijeme",
            "",
            "🔮 Napredne usporedbe",
            "   • Više od 2 laptopa, performance grafovi",
            "",
            "🔮 Deployment",
            "   • Docker kontejnerizacija, cloud hosting"
        ]
    )
    
    # Slide 18: Conclusion
    add_content_slide(
        "Zaključak",
        [
            "LaptopAI koristi modernu RAG arhitekturu koja kombinira:",
            "",
            "✅ Scraping za prikupljanje znanja",
            "",
            "✅ Vector bazu za brzo semantičko pretraživanje",
            "",
            "✅ LLM za generiranje strukturiranih insights",
            "",
            "✅ REST API za integraciju s frontend-om",
            "",
            "✅ React UI za interaktivnu usporedbu laptopa",
            "",
            "",
            "Sustav je dizajniran modularno → lagano proširivanje i održavanje"
        ]
    )
    
    # Slide 19: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    thanks_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(2))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Hvala na pažnji!"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(48)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = TITLE_COLOR
    thanks_para.alignment = PP_ALIGN.CENTER
    
    thanks_frame.add_paragraph()
    p1 = thanks_frame.paragraphs[1]
    p1.text = ""
    
    thanks_frame.add_paragraph()
    p2 = thanks_frame.paragraphs[2]
    p2.text = "Pitanja?"
    p2.font.size = Pt(32)
    p2.font.color.rgb = ACCENT_COLOR
    p2.alignment = PP_ALIGN.CENTER
    
    thanks_frame.add_paragraph()
    p3 = thanks_frame.paragraphs[3]
    p3.text = ""
    
    thanks_frame.add_paragraph()
    p4 = thanks_frame.paragraphs[4]
    p4.text = "Mario Radoš | PMF-ST | 2026"
    p4.font.size = Pt(16)
    p4.font.color.rgb = TEXT_COLOR
    p4.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('LaptopAI_Prezentacija.pptx')
    print("Prezentacija uspjesno kreirana: LaptopAI_Prezentacija.pptx")
    print(f"Ukupno slajdova: {len(prs.slides)}")

if __name__ == "__main__":
    create_presentation()
